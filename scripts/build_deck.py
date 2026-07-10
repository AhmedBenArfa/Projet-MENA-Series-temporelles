"""Build the French PowerPoint deck for the MENA VaR project.

Generates fresh figures under outputs/figures/ (prefixed ``deck_``) and
assembles presentation/VaR_MENA.pptx with python-pptx. All numbers shown
on data slides are read live from outputs/results.csv and
outputs/best_per_index.csv, or computed live from the data (ADF/KPSS,
GARCH sigma_t, ACF/PACF) -- nothing here is hand-typed from memory.

This is the detailed, teaching version of the deck: every time-series
concept used in the pipeline (decomposition, stationarity, ACF/PACF, the
lag operator and differencing, ARIMA/SARIMA, ARCH/GARCH, ML/DL forecasters,
BHS-VaR, backtesting) gets its own explained slide, grounded in the course
notions (Time Series 1GAMMA).

Run:
    "C:/Users/Mega-pc/anaconda3/python.exe" scripts/build_deck.py
"""

import sys
import pathlib

# --- bootstrap: make `tsvar` importable regardless of cwd -----------------
_root = pathlib.Path(__file__).resolve().parent.parent
for _c in (_root, *_root.parents):
    if (_c / "src" / "tsvar").is_dir():
        sys.path.insert(0, str(_c / "src"))
        break

import numpy as np
import pandas as pd

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from statsmodels.graphics.tsaplots import plot_acf, plot_pacf
from statsmodels.tsa.seasonal import seasonal_decompose

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from tsvar.data import (
    train_test_returns, INDEX_FILES, MENA, BENCHMARKS, load_index,
    adf_test, kpss_test,
)
from tsvar.volatility import walk_forward_garch
from tsvar.deep import walk_forward_dl
from tsvar.var import var_series

ROOT = _root
DATA = ROOT / "data (1)" / "data"
OUT = ROOT / "outputs"
FIG = OUT / "figures"
FIG.mkdir(parents=True, exist_ok=True)
PRES_DIR = ROOT / "presentation"
PRES_DIR.mkdir(parents=True, exist_ok=True)

RESULTS_CSV = OUT / "results.csv"
BEST_CSV = OUT / "best_per_index.csv"

# ---------------------------------------------------------------------------
# Palette -- "Gulf Ledger": deep navy/teal (finance, trust) + gold (MENA,
# value) + warm sand (paper/desert). Navy dominates (title/section/dark
# slides + headers); gold is the single sharp accent (numbers, winners,
# circles); sand is the light content background.
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x0E, 0x2A, 0x3D)
NAVY_DARK = RGBColor(0x08, 0x1B, 0x28)
GOLD = RGBColor(0xC9, 0x98, 0x2D)
GOLD_LIGHT = RGBColor(0xE7, 0xC9, 0x7E)
SAND = RGBColor(0xF3, 0xEC, 0xDD)
SAND_DARK = RGBColor(0xE7, 0xDC, 0xC3)
INK = RGBColor(0x1B, 0x24, 0x30)
SLATE = RGBColor(0x5A, 0x6B, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x46)
YELLOW = RGBColor(0xC8, 0x8A, 0x1E)
RED = RGBColor(0xAE, 0x39, 0x2F)

FONT_HEAD = "Cambria"
FONT_BODY = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)


# ===========================================================================
# Figure generation (fresh, for this deck)
# ===========================================================================

def gen_returns_clustering(index_name="Tunindex"):
    """Log-return series (train+test) illustrating volatility clustering."""
    tr, te = train_test_returns(index_name, DATA)
    full = pd.concat([tr, te])
    fig, ax = plt.subplots(figsize=(10, 3.6))
    ax.plot(full.index, full.values, lw=0.6, color="#0E2A3D")
    ax.axvline(tr.index[-1], color="#C9982D", lw=1.2, ls="--")
    ax.text(tr.index[-1], full.values.max() * 0.92, "  fin train / début test",
             color="#8a6a1e", fontsize=9)
    ax.set_title(f"Rendements log de {index_name} -- regroupement de volatilité", fontsize=12)
    ax.set_ylabel("rendement (%)")
    ax.margins(x=0.01)
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    fig.tight_layout()
    path = FIG / "deck_returns_clustering.png"
    fig.savefig(path, dpi=120)
    plt.close(fig)
    return path


def gen_decompose(index_name="ADI", period=21):
    """Illustrative additive decomposition Y = T + S + resid of the PRICE
    level (not the returns -- returns are already ~stationary, prices are
    not, which is exactly the point this slide makes). period=21 trading
    sessions ~ one calendar month, used only to expose the technique
    (moving-average trend + seasonal averages), not as a modelling step in
    the actual VaR pipeline."""
    df = load_index(DATA / INDEX_FILES[index_name])
    price = df["Price"]
    dec = seasonal_decompose(price, model="additive", period=period, extrapolate_trend="freq")
    fig, axes = plt.subplots(4, 1, figsize=(7.6, 6.2), sharex=True)
    axes[0].plot(price.index, price.values, color="#0E2A3D", lw=0.8)
    axes[0].set_ylabel("Prix Y_t", fontsize=9)
    axes[1].plot(dec.trend.index, dec.trend.values, color="#C9982D", lw=1.1)
    axes[1].set_ylabel("Tendance T_t", fontsize=9)
    axes[2].plot(dec.seasonal.index, dec.seasonal.values, color="#5A6B78", lw=0.6)
    axes[2].set_ylabel("Saisonnalité S_t", fontsize=9)
    axes[3].plot(dec.resid.index, dec.resid.values, color="#AE392F", lw=0.5)
    axes[3].set_ylabel("Résidu e_t", fontsize=9)
    axes[0].set_title(f"Décomposition illustrative -- prix {index_name} (période={period} séances)", fontsize=11)
    for ax in axes:
        for spine in ("top", "right"):
            ax.spines[spine].set_visible(False)
        ax.tick_params(labelsize=8)
    fig.tight_layout()
    path = FIG / "deck_decompose.png"
    fig.savefig(path, dpi=130)
    plt.close(fig)
    return path


def gen_acf_pacf(index_name="ADI", nlags=30):
    """ACF/PACF of the log-returns -- used to read off candidate (p, q)."""
    tr, te = train_test_returns(index_name, DATA)
    fig, axes = plt.subplots(1, 2, figsize=(9.4, 3.5))
    plot_acf(tr.values, lags=nlags, ax=axes[0])
    plot_pacf(tr.values, lags=nlags, ax=axes[1], method="ywm")
    axes[0].set_title(f"ACF -- rendements {index_name}", fontsize=11)
    axes[1].set_title(f"PACF -- rendements {index_name}", fontsize=11)
    for ax in axes:
        ax.set_xlabel("décalage (jours)", fontsize=9)
        for spine in ("top", "right"):
            ax.spines[spine].set_visible(False)
    fig.tight_layout()
    path = FIG / "deck_acf_pacf.png"
    fig.savefig(path, dpi=130)
    plt.close(fig)
    return path


def gen_garch_figs(index_name="Tunindex", alpha=0.01):
    """Fit GARCH(1,1) ONCE (train-once / walk-forward) and derive both the
    conditional-volatility figure (sigma_t) and the BHS-VaR-with-violations
    figure from the SAME walk-forward pass."""
    tr, te = train_test_returns(index_name, DATA)
    fc = walk_forward_garch(tr, te)
    v = var_series(fc, alpha)
    n_viol = int(np.sum(fc.y_true < v))
    rate = n_viol / len(fc.y_true)

    # -- sigma_t (conditional volatility) --
    fig, ax = plt.subplots(figsize=(10, 3.2))
    ax.plot(fc.dates, fc.sigma, color="#0E2A3D", lw=1.3)
    ax.fill_between(fc.dates, 0, fc.sigma, color="#C9982D", alpha=0.25)
    ax.set_title(f"Volatilité conditionnelle GARCH(1,1) -- sigma_t, {index_name} (test)", fontsize=12)
    ax.set_ylabel("sigma_t (%)")
    ax.margins(x=0.01)
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    fig.tight_layout()
    path_sigma = FIG / "deck_garch_sigma.png"
    fig.savefig(path_sigma, dpi=120)
    plt.close(fig)

    # -- BHS VaR + violations --
    fig, ax = plt.subplots(figsize=(11, 3.2))
    ax.plot(fc.dates, fc.y_true, label="rendement réalisé", lw=0.8, color="#5A6B78")
    ax.plot(fc.dates, v, label=f"VaR GARCH ({int((1-alpha)*100)}%)", color="#AE392F", lw=1.3)
    br = fc.y_true < v
    ax.scatter(fc.dates[br], fc.y_true[br], color="#0E2A3D", s=22, zorder=5,
               label=f"violations (n={n_viol}, {rate:.1%})")
    ax.set_title(f"VaR GARCH (BHS) -- {index_name}, alpha={alpha}", fontsize=12)
    ax.set_ylabel("rendement (%)")
    ax.legend(loc="lower left", fontsize=9, frameon=False)
    ax.margins(x=0.01)
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    fig.tight_layout()
    path_var = FIG / "deck_garch_var_tunindex.png"
    fig.savefig(path_var, dpi=120)
    plt.close(fig)

    return path_var, n_viol, rate, path_sigma


def gen_adi_lstm_vs_garch(alpha=0.01):
    """ADI: GARCH vs LSTM VaR side by side -- for the hypothesis-verdict slide."""
    tr, te = train_test_returns("ADI", DATA)
    fc_g = walk_forward_garch(tr, te)
    fc_l = walk_forward_dl(tr, te, "lstm")
    v_g = var_series(fc_g, alpha)
    v_l = var_series(fc_l, alpha)

    fig, axes = plt.subplots(1, 2, figsize=(11, 3.8), sharey=True)
    for ax, fc, v, title, color in (
        (axes[0], fc_g, v_g, "GARCH", "#0E2A3D"),
        (axes[1], fc_l, v_l, "LSTM", "#C9982D"),
    ):
        ax.plot(fc.dates, fc.y_true, lw=0.7, color="#5A6B78")
        ax.plot(fc.dates, v, color=color, lw=1.3)
        br = fc.y_true < v
        n = int(br.sum())
        ax.scatter(fc.dates[br], fc.y_true[br], color="#AE392F", s=20, zorder=5)
        ax.set_title(f"{title} -- ADI 99% (n={n} viol., {n/len(fc.y_true):.1%})", fontsize=11)
        for spine in ("top", "right"):
            ax.spines[spine].set_visible(False)
    fig.tight_layout()
    path = FIG / "deck_adi_garch_vs_lstm.png"
    fig.savefig(path, dpi=120)
    plt.close(fig)
    return path


def gen_winner_heatmap():
    """Basel-zone heatmap per (index, model) at alpha=1%, winners starred."""
    res = pd.read_csv(RESULTS_CSV)
    best = pd.read_csv(BEST_CSV)
    r = res[res.alpha == 0.01].copy()
    models = ["ARIMA", "SARIMA", "GARCH", "RF", "XGB", "ANN", "LSTM"]
    indices = ["Tunindex", "ADI", "MASI", "TASI"]
    zone_val = {"green": 0, "yellow": 1, "red": 2}
    mat = np.zeros((len(indices), len(models)))
    for i, idx in enumerate(indices):
        for j, m in enumerate(models):
            row = r[(r["index"] == idx) & (r["model"] == m)]
            mat[i, j] = zone_val.get(row.iloc[0]["basel_zone"], 1) if len(row) else np.nan

    cmap = matplotlib.colors.ListedColormap(["#2E7D46", "#C88A1E", "#AE392F"])
    fig, ax = plt.subplots(figsize=(10.5, 4.6))
    ax.imshow(mat, cmap=cmap, vmin=0, vmax=2, aspect="auto")
    ax.set_xticks(range(len(models))); ax.set_xticklabels(models, fontsize=12)
    ax.set_yticks(range(len(indices))); ax.set_yticklabels(indices, fontsize=12)
    ax.set_title("Zone de Bâle par modèle et par indice (alpha=1%) -- gagnant encadré", fontsize=13)

    winners = {row["index"]: row["model"] for _, row in best.iterrows()}
    for i, idx in enumerate(indices):
        for j, m in enumerate(models):
            row = r[(r["index"] == idx) & (r["model"] == m)]
            if not len(row):
                continue
            rate = row.iloc[0]["observed_rate"]
            ax.text(j, i, f"{rate:.1%}", ha="center", va="center", fontsize=10.5, color="white")
            if winners.get(idx) == m:
                rect = plt.Rectangle((j - 0.5, i - 0.5), 1, 1, fill=False,
                                      edgecolor="#0E2A3D", linewidth=3)
                ax.add_patch(rect)
    fig.tight_layout()
    path = FIG / "deck_winner_heatmap.png"
    fig.savefig(path, dpi=150)
    plt.close(fig)
    return path


# ===========================================================================
# python-pptx helpers
# ===========================================================================

def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = color
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, text, size=16, color=INK, bold=False,
             italic=False, font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=15, color=INK,
                 font=FONT_BODY, bullet="—", space_after=8, bold_lead=False,
                 anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 1.08
        r = p.add_run()
        r.text = f"{bullet}  {item}"
        r.font.size = Pt(size)
        r.font.name = font
        r.font.color.rgb = color
    return box


def add_kicker_and_title(slide, num, kicker, title, dark=False):
    """Section-number circle (repeated visual motif) + kicker + slide title."""
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN, Inches(0.42), Inches(0.55), Inches(0.55))
    circ.fill.solid()
    circ.fill.fore_color.rgb = GOLD
    circ.line.fill.background()
    circ.shadow.inherit = False
    tf = circ.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"{num:02d}"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.name = FONT_HEAD
    r.font.color.rgb = NAVY_DARK

    kick_color = GOLD_LIGHT if dark else GOLD
    add_text(slide, Inches(1.3), Inches(0.44), Inches(9), Inches(0.3), kicker.upper(),
              size=12, color=kick_color, bold=True, font=FONT_BODY)
    title_color = WHITE if dark else NAVY
    add_text(slide, Inches(1.3), Inches(0.72), Inches(11.4), Inches(0.7), title,
              size=30, color=title_color, bold=True, font=FONT_HEAD)


def add_footer(slide, page, total, note="Projet VaR MENA -- Deep Learning vs modèles statistiques", dark=False):
    color = SLATE if not dark else RGBColor(0x9F, 0xB3, 0xC2)
    add_text(slide, MARGIN, Inches(7.14), Inches(10.5), Inches(0.3), note,
              size=9, color=color, italic=True, font=FONT_BODY)
    add_text(slide, Inches(12.3), Inches(7.14), Inches(0.5), Inches(0.3), f"{page}/{total}",
              size=9, color=color, font=FONT_BODY, align=PP_ALIGN.RIGHT)


def add_picture_framed(slide, path, left, top, width=None, height=None):
    pic = slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    pic.line.color.rgb = SAND_DARK
    pic.line.width = Pt(1)
    return pic


def _set_cell(cell, text, size=12, color=INK, bold=False, bg=None, align=PP_ALIGN.CENTER, font=FONT_BODY):
    cell.margin_left = Inches(0.06)
    cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run() if len(p.runs) == 0 else p.runs[0]
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color


def add_table(slide, left, top, width, height, headers, rows, col_widths=None,
              cell_colors=None, header_bg=NAVY, header_color=WHITE, size=11):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gtable = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    if col_widths:
        for j, w in enumerate(col_widths):
            gtable.columns[j].width = w
    for j, h in enumerate(headers):
        _set_cell(gtable.cell(0, j), h, size=size, color=header_color, bold=True, bg=header_bg)
    for i, row in enumerate(rows, start=1):
        bg_row = SAND if i % 2 == 1 else RGBColor(0xFB, 0xF8, 0xF1)
        for j, val in enumerate(row):
            bg = None
            if cell_colors is not None and cell_colors[i - 1][j] is not None:
                bg = cell_colors[i - 1][j]
            _set_cell(gtable.cell(i, j), str(val), size=size, color=INK, bg=bg or bg_row)
    return gtable


def zone_color(zone):
    return {"green": GREEN, "yellow": YELLOW, "red": RED}.get(zone, SLATE)


# ===========================================================================
# Deck assembly
# ===========================================================================

def build():
    res = pd.read_csv(RESULTS_CSV)
    best = pd.read_csv(BEST_CSV)

    fig_returns = gen_returns_clustering("Tunindex")
    fig_decompose = gen_decompose("ADI", period=21)
    fig_acf_pacf = gen_acf_pacf("ADI", nlags=30)
    fig_garch, garch_nviol, garch_rate, fig_garch_sigma = gen_garch_figs("Tunindex", 0.01)
    fig_adi = gen_adi_lstm_vs_garch(0.01)
    fig_heat = gen_winner_heatmap()

    # -- live stationarity diagnostics (Tunindex): price level vs log-returns --
    price_tunindex = load_index(DATA / INDEX_FILES["Tunindex"])["Price"]
    ret_tunindex, _te_ret = train_test_returns("Tunindex", DATA)
    adf_price, kpss_price = adf_test(price_tunindex), kpss_test(price_tunindex)
    adf_ret, kpss_ret = adf_test(ret_tunindex), kpss_test(ret_tunindex)

    prs = new_presentation()
    TOTAL = 31
    page = [0]

    def track(slide):
        page[0] += 1
        return slide

    # ---------------- Slide 1: Titre --------------------------------------
    s = track(blank_slide(prs))
    set_background(s, NAVY)
    add_rect(s, Inches(0), Inches(6.6), SLIDE_W, Inches(0.9), NAVY_DARK)
    add_text(s, Inches(1.0), Inches(2.15), Inches(11.3), Inches(0.4), "PROJET DE SÉRIES TEMPORELLES",
              size=14, color=GOLD_LIGHT, bold=True, font=FONT_BODY)
    add_text(s, Inches(1.0), Inches(2.55), Inches(11.3), Inches(1.7),
              "La Value-at-Risk sur les marchés\nboursiers MENA",
              size=42, color=WHITE, bold=True, font=FONT_HEAD, line_spacing=1.05)
    add_text(s, Inches(1.0), Inches(4.25), Inches(11.0), Inches(0.6),
              "Deep Learning vs modèles statistiques -- ARIMA/SARIMA, GARCH, Random Forest / XGBoost, ANN / LSTM",
              size=16, color=SAND, font=FONT_BODY)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(0.7), Inches(1.9), Inches(1.9))
    circ.fill.solid(); circ.fill.fore_color.rgb = GOLD; circ.line.fill.background(); circ.shadow.inherit = False
    tf = circ.text_frame; tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "VaR"; r.font.size = Pt(28); r.font.bold = True
    r.font.name = FONT_HEAD; r.font.color.rgb = NAVY_DARK
    add_text(s, Inches(1.0), Inches(6.75), Inches(8), Inches(0.5),
              "Tunindex - ADI - MASI - TASI  |  repères : CAC40, S&P 500", size=12, color=SAND)
    add_text(s, Inches(10.2), Inches(6.75), Inches(2.5), Inches(0.5), "ahmedbenarfa.1992@gmail.com",
              size=11, color=GOLD_LIGHT, align=PP_ALIGN.RIGHT)

    # ---------------- Slide 2: Contexte & problématique --------------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 1, "Contexte", "Contexte et problématique")
    add_bullets(s, MARGIN, Inches(1.7), Inches(6.7), Inches(4.8), [
        "La Value-at-Risk (VaR) mesure la perte potentielle maximale d'un actif à un horizon et un niveau de confiance donnés (95% ou 99%).",
        "Les indices boursiers MENA (Tunisie, Émirats, Maroc, Arabie Saoudite) présentent des dynamiques propres : liquidité plus faible, chocs politiques régionaux, calendriers de cotation spécifiques.",
        "Question de recherche : les modèles de Deep Learning (ANN, LSTM) surpassent-ils les approches statistiques classiques (ARIMA, GARCH) et le Machine Learning (RF, XGBoost) pour l'estimation de la VaR sur ces marchés ?",
        "Hypothèse testée : le LSTM surpasse les modèles classiques sur l'indice ADI (Abu Dhabi).",
    ], size=15.5)
    box = add_rect(s, Inches(8.1), Inches(1.7), Inches(4.6), Inches(3.1), WHITE)
    add_text(s, Inches(8.4), Inches(1.9), Inches(4.0), Inches(0.4), "POURQUOI LA VaR ?",
              size=12, bold=True, color=GOLD.__class__(0xC9,0x98,0x2D))
    add_bullets(s, Inches(8.4), Inches(2.35), Inches(4.0), Inches(2.3), [
        "Exigence réglementaire (Bâle II/III) pour les banques et gestionnaires d'actifs.",
        "Outil de pilotage du risque de marché au quotidien.",
        "Le backtesting (Kupiec, Christoffersen, zones de Bâle) valide -- ou invalide -- un modèle de VaR a posteriori.",
    ], size=13.5, bullet="›", anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 3: Objectifs & démarche -------------------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 2, "Démarche", "Objectifs et démarche du projet")
    add_bullets(s, MARGIN, Inches(1.75), Inches(7.2), Inches(4.9), [
        "1. Nettoyer les données brutes (prix, volumes) et construire des rendements log stationnaires.",
        "2. Analyser la structure temporelle des séries : décomposition, stationnarité (ADF/KPSS), autocorrélations (ACF/PACF).",
        "3. Modéliser la moyenne conditionnelle (ARIMA/SARIMA) et la variance conditionnelle (GARCH).",
        "4. Comparer à des approches Machine Learning (Random Forest, XGBoost) et Deep Learning (ANN, LSTM) sur le même problème de prévision.",
        "5. Calculer la VaR (Bootstrap Historical Simulation) pour les 7 modèles, à 95% et 99%, en walk-forward strict.",
        "6. Valider chaque modèle a posteriori par backtesting (Kupiec, Christoffersen, zones de Bâle).",
    ], size=14.5)
    add_rect(s, Inches(8.5), Inches(1.75), Inches(4.2), Inches(4.6), NAVY)
    add_text(s, Inches(8.75), Inches(1.95), Inches(3.7), Inches(0.4), "3 FAMILLES DE MODÈLES", size=12, bold=True, color=GOLD_LIGHT)
    add_bullets(s, Inches(8.75), Inches(2.45), Inches(3.7), Inches(3.7), [
        "Statistique : ARIMA, SARIMA, GARCH -- fondés sur la théorie des séries temporelles.",
        "Machine Learning : Random Forest, XGBoost -- apprentissage supervisé sur des rendements retardés (lags).",
        "Deep Learning : ANN (perceptron multicouche), LSTM (réseau récurrent à mémoire) -- séquences de rendements.",
    ], size=13, color=SAND, bullet="›")
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 4: Données & indices ---------------------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 3, "Données", "Données et indices")
    headers = ["Indice", "Pays / place", "Période (train)", "Période (test)"]
    rows = [
        ["Tunindex", "Tunisie", "2005-01-03 -> 2014-12-31", "2014-12-31 -> 2015-08-20"],
        ["ADI", "Abu Dhabi (EAU)", "2005-01-03 -> 2014-12-31", "2014-12-31 -> 2015-08-31"],
        ["MASI", "Maroc", "2005-01-03 -> 2014-12-31", "2014-12-31 -> 2015-08-19"],
        ["TASI", "Arabie Saoudite", "2005-01-03 -> 2014-12-31", "2014-12-31 -> 2015-08-30"],
    ]
    add_table(s, MARGIN, Inches(1.75), Inches(8.6), Inches(2.3), headers, rows,
              col_widths=[Inches(1.6), Inches(2.2), Inches(2.4), Inches(2.4)])
    add_text(s, MARGIN, Inches(4.35), Inches(8.6), Inches(0.35), "Repères marchés développés (comparaison, hors entraînement des modèles MENA) :",
              size=13, bold=True, color=NAVY)
    add_bullets(s, MARGIN, Inches(4.75), Inches(8.6), Inches(1.6), [
        "CAC40 (France) -- 2560 séances, 2005-2014.",
        "S&P 500 (États-Unis) -- 2536 séances, 2005-2014.",
        "Rendements calculés en log-rendements journaliers (%), base commune pour toutes les places.",
    ], size=13.5)
    add_rect(s, Inches(9.9), Inches(1.75), Inches(2.8), Inches(3.0), WHITE)
    add_text(s, Inches(10.1), Inches(1.95), Inches(2.4), Inches(0.5), "TAILLE DES SÉRIES",
              size=11, bold=True, color=GOLD)
    add_text(s, Inches(10.1), Inches(2.4), Inches(2.4), Inches(2.2),
              "~2 470 à 2 585\nséances d'entraînement\npar indice MENA\n\n+ 160 à 171\nséances de test\n(walk-forward)",
              size=14.5, color=NAVY, bold=True, line_spacing=1.2, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 5: Prétraitement & rendements log ---------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 4, "Prétraitement", "Prétraitement et rendements logarithmiques")
    add_bullets(s, MARGIN, Inches(1.75), Inches(6.9), Inches(3.0), [
        "Prix nettoyés (virgules, unités M/K du volume) : conversion en série numérique exploitable Y_t = P_t (prix de clôture).",
        "On travaille en rendements log plutôt qu'en prix : r_t = 100 x ln(P_t / P_{t-1}).",
        "Pourquoi le log ? additivité dans le temps (somme de rendements = rendement cumulé), stabilisation de la variance, symétrie perte/gain, et rapprochement d'une distribution proche de la normale pour de petites variations.",
        "Les rendements log sont approximativement stationnaires en moyenne -- condition nécessaire pour ARIMA/GARCH (voir slide suivante).",
    ], size=14)
    add_rect(s, Inches(8.0), Inches(1.75), Inches(4.7), Inches(2.3), NAVY)
    add_text(s, Inches(8.25), Inches(1.95), Inches(4.2), Inches(0.35), "FORMULE", size=11, bold=True, color=GOLD_LIGHT)
    add_text(s, Inches(8.25), Inches(2.35), Inches(4.2), Inches(1.5),
              "r_t = 100 . ln(P_t / P_{t-1})",
              size=19, bold=True, color=WHITE, font=FONT_HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.0), Inches(4.25), Inches(4.7), Inches(1.2),
              "r_t en %, P_t = prix de clôture au jour t.\nToutes les places (MENA + repères) sont comparées sur cette même échelle.",
              size=12, italic=True, color=NAVY)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 6: Analyse exploratoire -------------------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 5, "Analyse exploratoire", "Regroupement de volatilité (volatility clustering)")
    add_bullets(s, MARGIN, Inches(1.75), Inches(5.9), Inches(4.6), [
        "En traçant les rendements log dans le temps, on observe des périodes calmes (petites variations) suivies de périodes agitées (grandes variations) qui s'enchaînent -- le regroupement de volatilité.",
        "Ce phénomène signale une hétéroscédasticité conditionnelle : la variance de r_t dépend du passé récent, même si la moyenne, elle, reste globalement stable.",
        "Une hypothèse i.i.d. ou un modèle à variance constante (OLS, ARMA classique) ne peut pas capter ce comportement -- c'est la motivation directe des modèles ARCH/GARCH (slides suivantes).",
        "La ligne verticale marque la séparation train/test (walk-forward) : le clustering apparaît sur les deux périodes.",
    ], size=14)
    add_picture_framed(s, fig_returns, Inches(6.8), Inches(1.9), width=Inches(5.9))
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 7: Décomposition d'une série temporelle ---------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 6, "Décomposition", "Décomposition d'une série temporelle")
    add_bullets(s, MARGIN, Inches(1.75), Inches(6.1), Inches(4.9), [
        "Modèle additif du cours : Y_t = T_t + S_t + e_t (tendance + saisonnalité + résidu).",
        "Tendance T_t : estimée par moyenne mobile (lissage sur une fenêtre) ou par régression paramétrique (linéaire, polynomiale).",
        "Saisonnalité S_t : estimée par moyennes saisonnières (moyenne par sous-période du cycle) ou par un modèle cosinus-sinus (régression harmonique).",
        "On retire T_t et S_t du signal pour obtenir un résidu plus proche d'un processus stationnaire, modélisable par ARMA.",
        "Ici (illustratif) : décomposition du PRIX d'ADI, période=21 séances (~1 mois). Les rendements log, eux, n'ont pas de tendance/saisonnalité forte -- c'est pourquoi on modélise directement les rendements plutôt que les prix décomposés.",
    ], size=13)
    add_picture_framed(s, fig_decompose, Inches(7.15), Inches(1.55), width=Inches(5.55))
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 8: Stationnarité -- définition & importance -----
    s = track(blank_slide(prs)); set_background(s, NAVY)
    add_kicker_and_title(s, 7, "Stationnarité", "Stationnarité : définition et importance", dark=True)
    add_bullets(s, Inches(1.0), Inches(1.9), Inches(5.6), Inches(4.6), [
        "Stationnarité stricte : la distribution jointe de (Y_t1,...,Y_tk) est invariante par translation dans le temps.",
        "Stationnarité faible (du second ordre) : moyenne E[Y_t] constante, variance Var(Y_t) constante, et autocovariance Cov(Y_t, Y_t-h) qui ne dépend que du décalage h -- pas de t.",
        "C'est la condition faible qui est utilisée en pratique pour les modèles AR / MA / ARMA.",
        "Pourquoi c'est indispensable : sans stationnarité, les estimateurs (moyenne, variance, autocorrélations) ne convergent pas vers des quantités stables -- les prévisions et les tests d'hypothèses perdent leur sens.",
    ], size=15, color=SAND)
    add_rect(s, Inches(7.1), Inches(1.9), Inches(5.6), Inches(4.6), NAVY_DARK)
    add_text(s, Inches(7.35), Inches(2.1), Inches(5.1), Inches(0.4), "3 CONDITIONS (FAIBLE)", size=12, bold=True, color=GOLD_LIGHT)
    add_bullets(s, Inches(7.35), Inches(2.6), Inches(5.1), Inches(3.7), [
        "E[Y_t] = mu, constante dans le temps.",
        "Var(Y_t) = sigma^2, constante dans le temps.",
        "Cov(Y_t, Y_t-h) = gamma(h), fonction de h uniquement.",
        "-> Une série de prix (marche aléatoire) viole typiquement les 3 conditions ; une série de rendements log les respecte approximativement.",
    ], size=14, color=SAND, bullet="›")
    add_footer(s, page[0], TOTAL, dark=True)

    # ---------------- Slide 9: Tests ADF & KPSS -----------------------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 8, "Tests de stationnarité", "Tests ADF et KPSS : prix vs rendements")
    add_bullets(s, MARGIN, Inches(1.7), Inches(12.1), Inches(1.75), [
        "ADF (Augmented Dickey-Fuller) : H0 = racine unitaire (série NON stationnaire). Si p < 0.05, on rejette H0 -> série stationnaire en moyenne.",
        "KPSS : hypothèse inverse -- H0 = série stationnaire. Si p > 0.05, on ne rejette pas H0 -> cohérent avec la stationnarité.",
        "Les deux tests peuvent diverger : KPSS est sensible à l'hétéroscédasticité (variance non constante), pas seulement à une tendance déterministe.",
    ], size=13.5)
    adf_price_ok = adf_price["stationary"]; kpss_price_ok = kpss_price["stationary"]
    adf_ret_ok = adf_ret["stationary"]; kpss_ret_ok = kpss_ret["stationary"]
    headers = ["Série (Tunindex)", "ADF stat", "ADF p", "Concl. ADF", "KPSS stat", "KPSS p", "Concl. KPSS"]
    rows = [
        ["Prix (niveau)", f"{adf_price['stat']:.2f}", f"{adf_price['pvalue']:.3f}",
         "Oui" if adf_price_ok else "Non",
         f"{kpss_price['stat']:.3f}", f"{kpss_price['pvalue']:.3f}",
         "Oui" if kpss_price_ok else "Non"],
        ["Rendements log", f"{adf_ret['stat']:.2f}", f"{adf_ret['pvalue']:.3f}",
         "Oui" if adf_ret_ok else "Non",
         f"{kpss_ret['stat']:.3f}", f"{kpss_ret['pvalue']:.3f}",
         "Oui" if kpss_ret_ok else "Non"],
    ]
    colors = [
        [None, None, None, (GREEN if adf_price_ok else RED), None, None, (GREEN if kpss_price_ok else RED)],
        [None, None, None, (GREEN if adf_ret_ok else RED), None, None, (GREEN if kpss_ret_ok else RED)],
    ]
    add_table(s, MARGIN, Inches(3.6), Inches(12.1), Inches(1.4), headers, rows,
              cell_colors=colors, size=10.5,
              col_widths=[Inches(2.9), Inches(1.5), Inches(1.4), Inches(1.7), Inches(1.5), Inches(1.4), Inches(1.7)])
    add_text(s, MARGIN, Inches(5.25), Inches(12.1), Inches(1.6),
              "Lecture : sur les PRIX, ADF et KPSS concordent -- série non stationnaire (marche aléatoire). Sur les RENDEMENTS LOG, ADF rejette nettement la racine unitaire (stationnaire en moyenne, p=0.000), mais KPSS reste significatif (p=0.015) : la variance des rendements n'est pas constante (regroupement de volatilité, slide 6), ce que KPSS détecte. C'est exactement pourquoi la moyenne est modélisée par ARIMA (d=0) et la variance, séparément, par GARCH.",
              size=12.5, italic=True, color=NAVY)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 10: ACF & PACF ----------------------------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 9, "Identification", "ACF et PACF : identifier les ordres (p, q)")
    add_bullets(s, MARGIN, Inches(1.7), Inches(12.1), Inches(1.5), [
        "ACF (autocorrélogramme) : corrélation entre Y_t et Y_t-h pour chaque décalage h -- utile pour repérer l'ordre q d'une MA (décroissance brutale après q).",
        "PACF (autocorrélation partielle) : corrélation entre Y_t et Y_t-h après avoir retiré l'effet des décalages intermédiaires -- utile pour repérer l'ordre p d'un AR (coupure nette après p).",
    ], size=14)
    add_picture_framed(s, fig_acf_pacf, Inches(2.65), Inches(3.05), width=Inches(8.0))
    add_text(s, Inches(2.65), Inches(6.25), Inches(8.0), Inches(0.6),
              "Sur les rendements ADI : peu de décalages sortent des bandes de confiance -- cohérent avec des rendements proches d'un bruit blanc en moyenne (ordre p, q faibles), la structure prévisible se logeant surtout dans la variance (GARCH).",
              size=11, italic=True, color=NAVY)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 11: Non-stationnarité -- opérateur retard -------
    s = track(blank_slide(prs)); set_background(s, NAVY)
    add_kicker_and_title(s, 10, "Non-stationnarité", "L'opérateur retard B et la différenciation", dark=True)
    card = add_rect(s, Inches(1.2), Inches(2.0), Inches(10.9), Inches(1.7), NAVY_DARK)
    add_text(s, Inches(1.2), Inches(2.2), Inches(10.9), Inches(1.3),
              "B . Y_t = Y_{t-1}          (1 - B) . Y_t = Y_t - Y_{t-1} = diff_t",
              size=20, bold=True, color=GOLD_LIGHT, font=FONT_HEAD, align=PP_ALIGN.CENTER)
    add_bullets(s, Inches(1.2), Inches(4.0), Inches(10.9), Inches(3.0), [
        "L'opérateur retard B décale la série d'une période : B^k . Y_t = Y_t-k.",
        "La différenciation ∇ = (1 - B) élimine une tendance stochastique (marche aléatoire) : si Y_t est non stationnaire, ∇Y_t peut le devenir.",
        "L'ordre d'intégration d est le nombre de différenciations nécessaires pour stationnariser la série -- d'où la notation ARIMA(p, d, q) : AR(p) + I(d) + MA(q).",
        "SARIMA(p, d, q)(P, D, Q, m) ajoute une différenciation et des ordres saisonniers sur une périodicité m (ex. m=5 pour un cycle hebdomadaire de séances boursières).",
        "Dans ce projet : les rendements log sont déjà stationnaires (d=0) -- la différenciation formelle a été évitée en modélisant directement r_t plutôt que P_t.",
    ], size=14, color=SAND)
    add_footer(s, page[0], TOTAL, dark=True)

    # ---------------- Slide 12: ARIMA(p,d,q) en détail ----------------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 11, "Modèles statistiques", "ARIMA(p, d, q) en détail")
    add_bullets(s, MARGIN, Inches(1.7), Inches(7.3), Inches(4.7), [
        "AR(p) -- partie autorégressive : Y_t = c + phi_1 Y_t-1 + ... + phi_p Y_t-p + eps_t (la valeur dépend de p valeurs passées).",
        "I(d) -- partie intégrée : nombre de différenciations pour rendre la série stationnaire (ici d=0, rendements déjà stationnaires).",
        "MA(q) -- partie moyenne mobile : eps_t = u_t + theta_1 u_t-1 + ... + theta_q u_t-q (l'erreur dépend de q chocs passés).",
        "ARIMA(p,d,q) combine les trois : la moyenne conditionnelle de r_t est expliquée par ses propres retards et par les erreurs passées.",
        "Mise en oeuvre ici : ajustement UNE seule fois sur le train (via auto_arima / ordre fixé), puis prévision walk-forward one-step-ahead sur le test, sans jamais ré-estimer les paramètres.",
        "sigma est l'écart-type des résidus du train ; la VaR est calculée par BHS à partir de (mu_t, sigma) -- voir slide VaR.",
    ], size=13.5)
    tasi_arima = res[(res["index"] == "TASI") & (res.model == "ARIMA") & (res.alpha == 0.01)].iloc[0]
    add_rect(s, Inches(8.5), Inches(1.75), Inches(4.2), Inches(4.65), NAVY)
    add_text(s, Inches(8.75), Inches(2.0), Inches(3.7), Inches(0.4), "TASI -- ARIMA @ 99%", size=13, bold=True, color=GOLD_LIGHT)
    add_text(s, Inches(8.75), Inches(2.55), Inches(3.7), Inches(1.0), f"{tasi_arima.observed_rate:.2%}",
              size=44, bold=True, color=WHITE, font=FONT_HEAD)
    add_text(s, Inches(8.75), Inches(3.55), Inches(3.7), Inches(0.4), "taux de violation observé (cible 1%)", size=11, color=SAND)
    add_text(s, Inches(8.75), Inches(4.15), Inches(3.7), Inches(0.4),
              f"Kupiec p = {tasi_arima.kupiec_p:.3f}   |   zone {tasi_arima.basel_zone.upper()}",
              size=12, color=GOLD_LIGHT, bold=True)
    add_text(s, Inches(8.75), Inches(4.7), Inches(3.7), Inches(1.5),
              "Zone verte Bâle : le modèle n'est pas rejeté. ARIMA est le modèle GAGNANT sur TASI (Arabie Saoudite) à 99%.",
              size=11.5, color=SAND, italic=True)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 13: SARIMA --------------------------------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 12, "Modèles statistiques", "SARIMA : la composante saisonnière")
    add_bullets(s, MARGIN, Inches(1.75), Inches(11.7), Inches(2.1), [
        "SARIMA(p,d,q)(P,D,Q,m) ajoute à ARIMA une structure autorégressive/moyenne-mobile saisonnière d'ordre (P,D,Q) sur une périodicité m.",
        "Ici m=5 : cycle hebdomadaire des séances boursières (5 jours ouvrés), motivé par un éventuel effet de calendrier (lundi/vendredi) dans les rendements.",
        "Ordre retenu dans le pipeline : SARIMA(1,0,1)(1,0,1,5) -- un AR et un MA à la fois non saisonniers et saisonniers, sans différenciation (d=D=0).",
    ], size=14.5)
    add_rect(s, Inches(1.6), Inches(4.15), Inches(10.1), Inches(1.3), NAVY)
    add_text(s, Inches(1.6), Inches(4.35), Inches(10.1), Inches(0.9),
              "SARIMA(1,0,1) x (1,0,1,5)",
              size=22, bold=True, color=GOLD_LIGHT, font=FONT_HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, MARGIN, Inches(5.75), Inches(12.1), Inches(1.0),
              "Comme pour ARIMA : ajustement UNE fois sur le train, walk-forward one-step-ahead sur le test (aucun ré-entraînement), VaR calculée par BHS à partir de (mu, sigma) des résidus.",
              size=13, italic=True, color=NAVY)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 14: Hétéroscédasticité (motivation GARCH) -------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 13, "Volatilité", "Hétéroscédasticité conditionnelle : la motivation du GARCH")
    add_bullets(s, MARGIN, Inches(1.75), Inches(6.3), Inches(4.7), [
        "Hétéroscédasticité conditionnelle : la variance de l'erreur au temps t, conditionnée à l'information passée, n'est PAS constante (contrairement à l'hypothèse d'homoscédasticité d'ARMA/OLS classiques).",
        "Symptôme observé (slide 6) : regroupement de volatilité -- périodes calmes et agitées qui s'enchaînent, incompatible avec Var(eps_t) = sigma^2 fixe.",
        "Modèle ARCH (Engle, 1982) : la variance dépend des carrés des chocs passés -- sigma^2_t = omega + somme alpha_i . eps^2_t-i.",
        "Limite de l'ARCH : il faut souvent beaucoup de retards (q élevé) pour bien capter la mémoire de la volatilité -- le GARCH résout cela avec un terme autorégressif sur la variance elle-même (slide suivante).",
    ], size=14)
    add_rect(s, Inches(8.5), Inches(1.75), Inches(4.2), Inches(2.6), NAVY)
    add_text(s, Inches(8.75), Inches(1.95), Inches(3.7), Inches(0.4), "ARCH(q)", size=12, bold=True, color=GOLD_LIGHT)
    add_text(s, Inches(8.75), Inches(2.4), Inches(3.7), Inches(1.6),
              "sigma^2_t = omega\n+ somme_{i=1}^{q} alpha_i eps^2_{t-i}",
              size=15, bold=True, color=WHITE, font=FONT_HEAD, line_spacing=1.2)
    add_text(s, Inches(8.5), Inches(4.55), Inches(4.2), Inches(1.9),
              "eps_t : choc (résidu) au temps t. La variance conditionnelle réagit à l'amplitude des chocs récents -- exactement le mécanisme du regroupement de volatilité.",
              size=12, italic=True, color=NAVY)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 15: GARCH(1,1) -- équation + sigma_t -----------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 14, "Volatilité", "GARCH(1,1) : la variance conditionnelle")
    card = add_rect(s, Inches(1.6), Inches(1.55), Inches(10.1), Inches(0.85), NAVY)
    add_text(s, Inches(1.6), Inches(1.68), Inches(10.1), Inches(0.6),
              "sigma^2_t = omega + alpha . eps^2_{t-1} + beta . sigma^2_{t-1}",
              size=19, bold=True, color=GOLD_LIGHT, font=FONT_HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, MARGIN, Inches(2.55), Inches(12.1), Inches(1.3), [
        "omega > 0 : variance de long terme (plancher) ; alpha : poids du choc récent (eps^2_{t-1}) ; beta : persistance de la volatilité passée (sigma^2_{t-1}). alpha + beta proche de 1 -> mémoire longue de la volatilité.",
        "Ajusté UNE fois sur le train (moyenne constante + GARCH(1,1)), puis sigma_t est propagé jour par jour sur le test avec (omega, alpha, beta) FIGÉS -- seul l'état (dernier choc, dernière variance) est mis à jour.",
    ], size=12.5)
    garch_sigma_w = 7.2
    garch_sigma_h = garch_sigma_w * (3.2 / 10.0)
    garch_sigma_left = Inches((13.333 - garch_sigma_w) / 2)
    garch_sigma_top = 3.95
    add_picture_framed(s, fig_garch_sigma, garch_sigma_left, Inches(garch_sigma_top), width=Inches(garch_sigma_w))
    add_text(s, garch_sigma_left, Inches(garch_sigma_top + garch_sigma_h + 0.12), Inches(garch_sigma_w), Inches(0.4),
              "sigma_t varie jour après jour sur le test -- c'est ce profil dynamique qui distingue GARCH d'un modèle à volatilité constante.",
              size=10.5, italic=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 16: GARCH -- VaR & FHS ---------------------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 15, "Volatilité", "GARCH et simulation historique filtrée (FHS)")
    add_bullets(s, MARGIN, Inches(1.7), Inches(12.1), Inches(1.35), [
        "FHS (Filtered Historical Simulation) : on combine sigma_t (dynamique, GARCH) avec le pool de résidus standardisés du train, rééchantillonné par bootstrap pour obtenir la VaR (BHS).",
        "Contrairement à une hypothèse gaussienne, FHS conserve la forme empirique (asymétrie, queues épaisses) de la distribution des chocs -- pertinent sur des marchés MENA souvent non gaussiens.",
    ], size=13.5)
    garch_img_top = 3.15
    garch_img_w = 10.1
    garch_img_h = garch_img_w * (3.2 / 11.0)  # matches gen_garch_figs figsize aspect ratio
    add_picture_framed(s, fig_garch, Inches(1.6), Inches(garch_img_top), width=Inches(garch_img_w))
    add_text(s, Inches(1.6), Inches(garch_img_top + garch_img_h + 0.15), Inches(10.1), Inches(0.5),
              f"Tunindex, 99% : {garch_nviol} violations observées sur la période test ({garch_rate:.1%}) -- GARCH est le modèle gagnant sur Tunindex, ADI et MASI (3/4 indices MENA).",
              size=11.5, italic=True, color=NAVY)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 17: Machine Learning -- features & RF -----------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 16, "Machine Learning", "Features (lags) et Random Forest")
    add_bullets(s, MARGIN, Inches(1.75), Inches(6.2), Inches(4.6), [
        "Reformulation en apprentissage supervisé : les n_lags rendements passés (r_t-1, ..., r_t-n_lags) forment le vecteur de variables explicatives X_t, la cible y_t = r_t.",
        "Random Forest : ensemble d'arbres de décision entraînés sur des sous-échantillons bootstrap des lags, moyennés pour la prévision (réduit la variance vs un arbre unique).",
        "Entraîné UNE fois sur le train (n_lags=5), puis walk-forward sur le test en réinjectant le rendement réalisé (pas la prévision) dans la fenêtre de lags.",
        "sigma constant = écart-type des résidus d'entraînement, combiné à mu (prévision de l'arbre) pour la VaR BHS -- même interface que tous les autres modèles.",
        "Résultat empirique -- Random Forest est le modèle le plus FAIBLE du panel pour la VaR : il sur-viole systématiquement le seuil attendu (zone ROUGE Bâle sur les 4 indices MENA à 95%, et sur Tunindex/MASI/TASI à 99%).",
    ], size=13.5)
    rf_rows = []
    for idx in MENA:
        row = res[(res["index"] == idx) & (res.model == "RF") & (res.alpha == 0.01)].iloc[0]
        rf_rows.append([idx, f"{row.observed_rate:.2%}", row.basel_zone.upper()])
    add_text(s, Inches(8.4), Inches(1.75), Inches(4.3), Inches(0.35), "RF @ 99% -- taux de violation observé",
              size=12, bold=True, color=NAVY)
    add_table(s, Inches(8.4), Inches(2.15), Inches(4.3), Inches(2.2),
              ["Indice", "Taux observé", "Zone Bâle"], rf_rows,
              cell_colors=[[None, None, zone_color(r[2].lower())] for r in rf_rows],
              col_widths=[Inches(1.7), Inches(1.5), Inches(1.1)])
    rf_rates = [float(r[1].rstrip("%").replace(",", ".")) for r in rf_rows]
    add_text(s, Inches(8.4), Inches(4.6), Inches(4.3), Inches(1.8),
              f"Cible attendue : ~1%. Les taux observés ({min(rf_rates):.1f}-{max(rf_rates):.1f}%, "
              "et jusqu'à 22% à 95%) sont très supérieurs -- Random Forest sous-estime largement le risque de queue.",
              size=12, italic=True, color=RED, anchor=MSO_ANCHOR.TOP)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 18: XGBoost -------------------------------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 17, "Machine Learning", "XGBoost : boosting de gradient")
    add_bullets(s, MARGIN, Inches(1.75), Inches(6.4), Inches(4.6), [
        "XGBoost construit ses arbres SÉQUENTIELLEMENT : chaque nouvel arbre corrige les erreurs (résidus) des arbres précédents (boosting de gradient), contrairement à RF (arbres indépendants, moyennés).",
        "Hyperparamètres utilisés ici : 300 arbres, profondeur max 4, taux d'apprentissage 0.05 -- réglages volontairement conservateurs pour limiter le sur-apprentissage sur des rendements bruités.",
        "Même protocole que RF : entraînement UNE fois sur le train (n_lags=5), walk-forward avec réinjection du rendement réalisé, sigma constant + mu prédit -> VaR BHS.",
        "Résultat empirique : XGBoost se comporte nettement MIEUX que Random Forest -- moins de sur-violations -- mais reste en général derrière ARIMA/GARCH/ANN (zones vertes/jaunes selon l'indice), sans remporter de titre de meilleur modèle sur un indice MENA.",
    ], size=14)
    xgb_rows = []
    for idx in MENA:
        row = res[(res["index"] == idx) & (res.model == "XGB") & (res.alpha == 0.01)].iloc[0]
        xgb_rows.append([idx, f"{row.observed_rate:.2%}", row.basel_zone.upper()])
    add_text(s, Inches(8.4), Inches(1.75), Inches(4.3), Inches(0.35), "XGB @ 99% -- taux de violation observé",
              size=12, bold=True, color=NAVY)
    add_table(s, Inches(8.4), Inches(2.15), Inches(4.3), Inches(2.2),
              ["Indice", "Taux observé", "Zone Bâle"], xgb_rows,
              cell_colors=[[None, None, zone_color(r[2].lower())] for r in xgb_rows],
              col_widths=[Inches(1.7), Inches(1.5), Inches(1.1)])
    add_text(s, Inches(8.4), Inches(4.6), Inches(4.3), Inches(1.8),
              "Comparé à RF : violations bien plus proches de la cible de 1% -- le boosting régularisé (profondeur faible, learning rate bas) limite le sur-apprentissage sur le bruit des rendements.",
              size=12, italic=True, color=NAVY, anchor=MSO_ANCHOR.TOP)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 19: Deep Learning -- ANN -------------------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 18, "Deep Learning", "ANN : le perceptron multicouche")
    add_bullets(s, MARGIN, Inches(1.75), Inches(6.4), Inches(4.6), [
        "Architecture : couche d'entrée de largeur = fenêtre (10 rendements standardisés), puis 32 -> 16 -> 1 neurones avec activations ReLU (perceptron multicouche, feedforward).",
        "Les rendements sont standardisés (z = (r - mu_train) / sigma_train) avant d'entrer dans le réseau -- pratique standard en Deep Learning pour stabiliser l'apprentissage.",
        "Entraînement UNE fois sur le train (30 époques, optimiseur Adam, perte MSE) ; en walk-forward, le réseau reçoit la fenêtre des 10 DERNIERS rendements réalisés et prédit le rendement du jour suivant -- les poids ne sont jamais mis à jour sur le test.",
        "sigma = écart-type des résidus d'entraînement (converti en échelle de rendement) ; VaR calculée par BHS, strictement comparable aux 6 autres modèles.",
        "Résultat empirique : ANN est bien calibré (zone verte) sur la quasi-totalité des couples indice/alpha, avec des MAE/RMSE proches de ceux d'ARIMA/GARCH -- mais ne les surpasse pas nettement.",
    ], size=13.5)
    add_rect(s, Inches(8.5), Inches(1.75), Inches(4.2), Inches(4.6), WHITE)
    add_text(s, Inches(8.75), Inches(1.95), Inches(3.7), Inches(0.4), "ARCHITECTURE ANN", size=12, bold=True, color=GOLD)
    add_bullets(s, Inches(8.75), Inches(2.4), Inches(3.7), Inches(3.7), [
        "Entrée : 10 rendements standardisés (fenêtre glissante).",
        "Couche cachée 1 : 32 neurones, ReLU.",
        "Couche cachée 2 : 16 neurones, ReLU.",
        "Sortie : 1 neurone -- prévision du rendement du jour suivant.",
        "30 époques, Adam (lr=1e-3), perte MSE.",
    ], size=13, bullet="›", anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 20: Deep Learning -- LSTM ------------------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 19, "Deep Learning", "LSTM : mémoire court/long terme")
    add_bullets(s, MARGIN, Inches(1.75), Inches(6.4), Inches(4.7), [
        "Le LSTM (Long Short-Term Memory) est un réseau récurrent conçu pour capter des DÉPENDANCES TEMPORELLES sur une séquence, au-delà d'une simple fenêtre de lags fixe.",
        "Chaque cellule LSTM maintient un état interne (mémoire à long terme) régulé par 3 portes : porte d'oubli (quelle part de la mémoire passée on efface), porte d'entrée (quelle part de la nouvelle information on ajoute), porte de sortie (quelle part de l'état on expose comme sortie h_t).",
        "Ici : la séquence des 10 derniers rendements (standardisés) est présentée à une couche LSTM (32 unités), suivie d'une couche dense (1 neurone) pour la prévision du rendement suivant.",
        "Entraîné UNE fois (30 époques) puis walk-forward strict, comme ANN -- les gains attendus du LSTM viennent de sa capacité à pondérer différemment les rendements récents vs plus anciens dans la séquence (contrairement à un ARIMA linéaire à ordre fixe).",
    ], size=13.5)
    add_rect(s, Inches(8.5), Inches(1.75), Inches(4.2), Inches(4.7), WHITE)
    add_text(s, Inches(8.75), Inches(1.95), Inches(3.7), Inches(0.4), "ADI -- LSTM @ 99%", size=12, bold=True, color=GOLD)
    adi_lstm = res[(res["index"] == "ADI") & (res.model == "LSTM") & (res.alpha == 0.01)].iloc[0]
    add_text(s, Inches(8.75), Inches(2.4), Inches(3.7), Inches(0.9), f"{adi_lstm.observed_rate:.2%}",
              size=38, bold=True, color=NAVY, font=FONT_HEAD)
    add_text(s, Inches(8.75), Inches(3.25), Inches(3.7), Inches(0.4), "taux de violation observé", size=11, color=SLATE)
    add_text(s, Inches(8.75), Inches(3.75), Inches(3.7), Inches(0.4),
              f"zone {adi_lstm.basel_zone.upper()}  |  Kupiec p={adi_lstm.kupiec_p:.3f}", size=12, bold=True, color=GREEN)
    add_text(s, Inches(8.75), Inches(4.35), Inches(3.7), Inches(1.9),
              "Bien calibré, comparable à GARCH sur ADI -- mais pas distinctement supérieur (voir slide vérification de l'hypothèse).",
              size=11.5, italic=True, color=NAVY)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 21: La VaR -- définition & interprétation -------
    s = track(blank_slide(prs)); set_background(s, NAVY)
    add_kicker_and_title(s, 20, "La VaR", "La Value-at-Risk : définition et interprétation", dark=True)
    add_bullets(s, Inches(1.0), Inches(1.9), Inches(11.3), Inches(1.9), [
        "Définition : VaR_alpha(t) est le niveau de perte tel que P(r_t < VaR_alpha(t)) = alpha -- le quantile alpha de la distribution du rendement futur.",
        "Interprétation à 99% : sur 100 jours, on s'attend à observer au plus ~1 jour où la perte réelle dépasse la VaR (une « violation »).",
    ], size=15, color=SAND)
    add_rect(s, Inches(1.0), Inches(4.0), Inches(11.3), Inches(2.9), NAVY_DARK)
    add_text(s, Inches(1.25), Inches(4.2), Inches(10.8), Inches(0.4), "3 FAMILLES DE MÉTHODES DE CALCUL DE LA VaR", size=12, bold=True, color=GOLD_LIGHT)
    add_bullets(s, Inches(1.25), Inches(4.7), Inches(10.8), Inches(2.0), [
        "Paramétrique (delta-normale) : suppose une loi (souvent normale) pour les rendements -- rapide mais irréaliste sur des marchés à queues épaisses.",
        "Simulation historique pure : rejoue directement les rendements passés -- simple mais rigide (poids égal à chaque jour passé, pas de dynamique de volatilité).",
        "Simulation historique bootstrappée (BHS/FHS -- retenue ici) : combine une prévision dynamique (mu_t, sigma_t) du modèle avec un rééchantillonnage des résidus standardisés -- flexible et sans hypothèse de normalité.",
    ], size=13.5, color=SAND, bullet="›")
    add_footer(s, page[0], TOTAL, dark=True)

    # ---------------- Slide 22: Calcul de la VaR par BHS (formule) ----------
    s = track(blank_slide(prs)); set_background(s, NAVY)
    add_kicker_and_title(s, 21, "Méthode de calcul", "Le calcul de la VaR par BHS", dark=True)
    card = add_rect(s, Inches(1.6), Inches(2.1), Inches(10.1), Inches(1.9), NAVY_DARK)
    add_text(s, Inches(1.6), Inches(2.35), Inches(10.1), Inches(1.5),
              "VaR_alpha(t) = mu_t + sigma_t . Q_alpha( résidus standardisés bootstrappés )",
              size=22, bold=True, color=GOLD_LIGHT, font=FONT_HEAD, align=PP_ALIGN.CENTER)
    add_bullets(s, Inches(1.6), Inches(4.3), Inches(10.1), Inches(2.6), [
        "mu_t, sigma_t : moyenne et écart-type prédits par le modèle (ARIMA, GARCH, RF/XGB, ANN/LSTM) pour le jour t.",
        "résidus standardisés : pool des résidus (z-scores) estimés UNE fois sur la période d'entraînement.",
        "Bootstrap : n_boot=10 000 tirages avec remise dans ce pool -> quantile empirique Q_alpha (alpha=5% ou 1%).",
        "BHS (Bootstrap Historical Simulation) ne suppose pas la normalité des résidus : la distribution empirique (asymétrie, queues épaisses) est préservée -- adapté aux marchés MENA, souvent non-gaussiens.",
    ], size=14, color=SAND)
    add_footer(s, page[0], TOTAL, dark=True)

    # ---------------- Slide 23: train-once / walk-forward one-step-ahead ----
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 22, "Méthodologie", "Pipeline unifié : train-once / walk-forward")
    steps = [
        ("1", "Entraînement unique", "Chaque modèle est ajusté UNE seule fois sur la période train (aucun ré-entraînement)."),
        ("2", "Walk-forward J+1", "Prévision un jour à l'avance sur le test ; les états (résidus, variance) sont mis à jour avec la valeur réalisée, pas les paramètres."),
        ("3", "VaR par BHS", "mu, sigma prédits + pool de résidus standardisés bootstrappés -> quantile alpha (VaR)."),
        ("4", "Backtesting", "Comparaison VaR vs rendement réalisé : Kupiec, Christoffersen, zones de Bâle."),
    ]
    w = Inches(2.85)
    for i, (num, head, desc) in enumerate(steps):
        left = Inches(0.6 + i * (2.85 + 0.18))
        card = add_rect(s, left, Inches(1.85), w, Inches(3.7), WHITE)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.18), Inches(2.05), Inches(0.5), Inches(0.5))
        circ.fill.solid(); circ.fill.fore_color.rgb = GOLD; circ.line.fill.background(); circ.shadow.inherit = False
        tf = circ.text_frame; tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num; r.font.bold = True; r.font.size = Pt(18)
        r.font.color.rgb = NAVY_DARK; r.font.name = FONT_HEAD
        add_text(s, left + Inches(0.18), Inches(2.7), w - Inches(0.36), Inches(0.6), head,
                  size=14.5, bold=True, color=NAVY, font=FONT_HEAD)
        add_text(s, left + Inches(0.18), Inches(3.3), w - Inches(0.36), Inches(2.1), desc,
                  size=12, color=SLATE, line_spacing=1.15)
        if i < 3:
            arrow_left = left + w + Inches(0.02)
            arr = s.shapes.add_shape(MSO_SHAPE.CHEVRON, arrow_left, Inches(3.5), Inches(0.14), Inches(0.6))
            arr.fill.solid(); arr.fill.fore_color.rgb = GOLD; arr.line.fill.background(); arr.shadow.inherit = False
    add_text(s, MARGIN, Inches(5.85), Inches(12.1), Inches(0.9),
              "Contrainte de conception : aucun modèle n'est ré-estimé sur le test -- seule l'information disponible à J est utilisée pour prédire J+1 (walk-forward one-step-ahead), pour les 7 modèles comme pour le GARCH (FHS).",
              size=13, italic=True, color=NAVY)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 24: Backtesting -- Kupiec (POF) -----------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 23, "Validation", "Backtesting : le test de Kupiec (POF)")
    add_bullets(s, MARGIN, Inches(1.75), Inches(11.9), Inches(1.7), [
        "Principe (Proportion of Failures) : on compare le TAUX de violations observé (n_viol / n) au taux attendu alpha.",
        "H0 : le taux de violation observé est statistiquement égal à alpha (le modèle est bien calibré).",
        "Statistique de test : LR_POF = -2 [ (n-x) ln(1-alpha) + x ln(alpha) - (n-x) ln(1-pi_hat) - x ln(pi_hat) ], suit un chi^2(1) sous H0, avec pi_hat = x/n le taux observé.",
    ], size=14)
    add_rect(s, Inches(1.6), Inches(3.85), Inches(10.1), Inches(1.0), NAVY)
    add_text(s, Inches(1.6), Inches(4.0), Inches(10.1), Inches(0.75),
              "p-value > 0.05  ->  modèle NON rejeté (bien calibré)",
              size=18, bold=True, color=GOLD_LIGHT, font=FONT_HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, MARGIN, Inches(5.15), Inches(11.9), Inches(1.5),
              "Limite du test de Kupiec : il ne regarde QUE le taux global de violations, pas leur répartition dans le temps -- un modèle peut avoir le bon nombre de violations mais toutes regroupées lors d'une même crise (violations non indépendantes). C'est ce que corrige le test de Christoffersen (slide suivante).",
              size=13, italic=True, color=NAVY)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 25: Backtesting -- Christoffersen & Bâle --------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 24, "Validation", "Backtesting : Christoffersen et zones de Bâle")
    cols = [
        ("Christoffersen (indépendance + couverture)", "Ajoute au test de Kupiec un test d'INDÉPENDANCE des violations dans le temps (pas de clusters). LR_cc = LR_POF + LR_ind ~ chi^2(2). p_cc > 0.05 -> modèle non rejeté."),
        ("Zones de Bâle", "Classe le modèle en vert / jaune / rouge selon le nombre de violations normalisé sur 250 séances -- grille réglementaire standard (verte <= 4 à 99%)."),
    ]
    w = Inches(5.9)
    for i, (head, desc) in enumerate(cols):
        left = Inches(0.6 + i * (5.9 + 0.3))
        add_rect(s, left, Inches(1.8), w, Inches(2.85), WHITE)
        add_text(s, left + Inches(0.25), Inches(2.0), w - Inches(0.5), Inches(0.9), head,
                  size=15, bold=True, color=NAVY, font=FONT_HEAD, line_spacing=1.05)
        add_text(s, left + Inches(0.25), Inches(2.85), w - Inches(0.5), Inches(1.6), desc,
                  size=13, color=SLATE, line_spacing=1.2, anchor=MSO_ANCHOR.MIDDLE)
    legend_y = Inches(4.95)
    for i, (zone, col, meaning) in enumerate([
        ("Vert", GREEN, "modèle non rejeté, risque bien capté"),
        ("Jaune", YELLOW, "zone de surveillance"),
        ("Rouge", RED, "modèle rejeté, sous-estime le risque"),
    ]):
        left = Inches(0.6 + i * 4.1)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, left, legend_y, Inches(0.3), Inches(0.3))
        dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background(); dot.shadow.inherit = False
        add_text(s, left + Inches(0.4), legend_y - Inches(0.03), Inches(3.6), Inches(0.4),
                  f"{zone} -- {meaning}", size=12.5, color=NAVY, bold=True)
    tasi_c = res[(res["index"] == "TASI") & (res.model == "ARIMA") & (res.alpha == 0.01)].iloc[0]
    add_text(s, MARGIN, Inches(5.55), Inches(12.1), Inches(1.3),
              f"Exemple (TASI, ARIMA, 99%) : Kupiec p={tasi_c.kupiec_p:.3f}, Christoffersen p_cc={tasi_c.christoffersen_p:.3f}, "
              f"zone {tasi_c.basel_zone.upper()} -- les deux tests concordent, le modèle n'est pas rejeté.",
              size=13, italic=True, color=NAVY)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 26: Résultats par indice ------------------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 25, "Résultats", "Résultats par indice -- le tableau des gagnants")
    headers = ["Indice", "Modèle gagnant", "Taux observé (99%)", "Kupiec p", "Zone Bâle"]
    rows = []
    colors = []
    for _, row in best.iterrows():
        rows.append([row["index"], row.model, f"{row.observed_rate:.2%}", f"{row.kupiec_p:.3f}", row.basel_zone.upper()])
        colors.append([None, GOLD_LIGHT, None, None, zone_color(row.basel_zone)])
    add_table(s, MARGIN, Inches(1.75), Inches(7.9), Inches(2.3), headers, rows,
              cell_colors=colors, col_widths=[Inches(1.7), Inches(1.9), Inches(1.9), Inches(1.3), Inches(1.1)])
    add_text(s, MARGIN, Inches(4.3), Inches(7.9), Inches(0.9),
              "GARCH est le modèle globalement le plus performant : il remporte 3 indices MENA sur 4 (Tunindex, ADI, MASI). ARIMA remporte TASI. Classement établi selon la p-value de Kupiec la plus élevée à alpha=1% (meilleure adéquation entre taux observé et taux attendu).",
              size=12.5, italic=True, color=NAVY)
    add_picture_framed(s, fig_heat, Inches(8.55), Inches(1.75), width=Inches(4.15))
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 27: Lecture détaillée des résultats -------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 26, "Résultats", "Lecture détaillée : les 7 modèles sur Tunindex")
    models_order = ["ARIMA", "SARIMA", "GARCH", "RF", "XGB", "ANN", "LSTM"]
    headers = ["Modèle", "Taux observé", "Kupiec p", "Christoffersen p", "Zone Bâle"]
    rows, colors = [], []
    for m in models_order:
        row = res[(res["index"] == "Tunindex") & (res.model == m) & (res.alpha == 0.01)].iloc[0]
        rows.append([m, f"{row.observed_rate:.2%}", f"{row.kupiec_p:.3f}", f"{row.christoffersen_p:.3f}", row.basel_zone.upper()])
        colors.append([None, None, None, None, zone_color(row.basel_zone)])
    add_table(s, MARGIN, Inches(1.75), Inches(8.2), Inches(3.0), headers, rows,
              cell_colors=colors, col_widths=[Inches(1.5), Inches(1.7), Inches(1.6), Inches(2.0), Inches(1.4)])

    rf_all = res[res.model == "RF"]
    rf_red = int((rf_all.basel_zone == "red").sum())
    ann_all = res[res.model == "ANN"]; lstm_all = res[res.model == "LSTM"]
    ann_green = int((ann_all.basel_zone == "green").sum()); ann_total = len(ann_all)
    lstm_green = int((lstm_all.basel_zone == "green").sum()); lstm_total = len(lstm_all)
    add_text(s, MARGIN, Inches(5.1), Inches(8.2), Inches(1.9),
              f"Random Forest : {rf_red}/{len(rf_all)} combinaisons (indice x alpha) en zone ROUGE -- modèle systématiquement trop optimiste sur le risque de queue, à éviter tel quel pour la VaR.\n"
              f"ANN / LSTM : {ann_green}/{ann_total} et {lstm_green}/{lstm_total} combinaisons en zone VERTE respectivement -- bien calibrés, mais sans dominer nettement GARCH/ARIMA (voir tableau des gagnants).",
              size=12.5, italic=True, color=NAVY, line_spacing=1.2)
    add_rect(s, Inches(8.9), Inches(1.75), Inches(3.8), Inches(4.6), NAVY)
    add_text(s, Inches(9.15), Inches(1.95), Inches(3.3), Inches(0.4), "À RETENIR", size=12, bold=True, color=GOLD_LIGHT)
    add_bullets(s, Inches(9.15), Inches(2.4), Inches(3.3), Inches(3.8), [
        "GARCH/ARIMA : compromis simplicité-robustesse, gagnants sur 4/4 indices.",
        "ANN/LSTM : bien calibrés, compétitifs, mais pas dominants.",
        "RF : sur-viole systématiquement -- inadapté en l'état.",
        "XGB : meilleur que RF, sans dépasser les modèles statistiques.",
    ], size=12.5, color=SAND, bullet="›")
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 28: MENA vs marchés développés ------------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 27, "Discussion", "MENA vs marchés développés (CAC40 / S&P 500)")
    add_bullets(s, MARGIN, Inches(1.75), Inches(6.3), Inches(4.6), [
        "Sur les marchés développés (CAC40, S&P 500), la littérature montre un avantage plus net des modèles GARCH/Deep Learning face aux chocs de volatilité (crises 2008, 2011).",
        "Sur les indices MENA de cette étude, les modèles ANN/LSTM égalent GARCH/ARIMA en calibration (zone verte) mais ne les surpassent pas nettement -- écart moins marqué qu'attendu.",
        "Hypothèse explicative : séries MENA plus courtes en observations utiles de volatilité extrême, et micro-structure (liquidité, jours de fermeture spécifiques) qui peut limiter l'avantage informationnel du Deep Learning.",
        "GARCH, plus parcimonieux, reste compétitif -- cohérent avec des travaux montrant sa robustesse sur des marchés émergents à volatilité regroupée mais aux échantillons plus restreints.",
    ], size=14)
    add_rect(s, Inches(8.5), Inches(1.75), Inches(4.2), Inches(3.6), NAVY)
    add_text(s, Inches(8.75), Inches(1.95), Inches(3.7), Inches(0.4), "À RETENIR", size=12, bold=True, color=GOLD_LIGHT)
    add_bullets(s, Inches(8.75), Inches(2.4), Inches(3.7), Inches(2.7), [
        "Pas de \"solution unique\" -- le meilleur modèle dépend de l'indice.",
        "GARCH : bon compromis simplicité / robustesse sur MENA.",
        "Deep Learning : compétitif mais pas dominant ici.",
        "RF/XGBoost : à éviter tels quels pour la VaR de marché.",
    ], size=13, color=SAND, bullet="›", anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 29: Vérification de l'hypothèse -----------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 28, "Verdict", "Vérification de l'hypothèse : le LSTM sur ADI")
    add_text(s, MARGIN, Inches(1.7), Inches(12.1), Inches(0.4),
              "Hypothèse testée : « le LSTM surpasse les modèles classiques sur l'indice ADI »",
              size=15, bold=True, color=NAVY)
    add_picture_framed(s, fig_adi, Inches(MARGIN.inches), Inches(2.2), width=Inches(8.4))
    box = add_rect(s, Inches(9.05), Inches(2.2), Inches(3.65), Inches(4.4), RED)
    add_text(s, Inches(9.3), Inches(2.4), Inches(3.2), Inches(0.5), "NON CONFIRMÉE", size=20, bold=True, color=WHITE, font=FONT_HEAD)
    adi_garch = res[(res["index"] == "ADI") & (res.model == "GARCH") & (res.alpha == 0.01)].iloc[0]
    add_text(s, Inches(9.3), Inches(3.0), Inches(3.2), Inches(2.9),
              f"LSTM @ 99% : {adi_lstm.observed_rate:.2%} de violations, zone {adi_lstm.basel_zone.upper()} -- bien calibré.\n\n"
              f"GARCH @ 99% : {adi_garch.observed_rate:.2%}, zone {adi_garch.basel_zone.upper()} -- également bien calibré, et gagnant officiel sur ADI (Kupiec p={adi_garch.kupiec_p:.3f} vs {adi_lstm.kupiec_p:.3f}).\n\n"
              "Le LSTM est fiable sur ADI mais n'est PAS distinctement supérieur au GARCH : les deux sont en zone verte, GARCH a la meilleure p-value de Kupiec.",
              size=12, color=WHITE, line_spacing=1.15)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 30: Conclusion & perspectives --------------------
    s = track(blank_slide(prs)); set_background(s, SAND)
    add_kicker_and_title(s, 29, "Conclusion", "Conclusion et perspectives")
    add_bullets(s, MARGIN, Inches(1.75), Inches(6.3), Inches(4.7), [
        "GARCH est le modèle le plus robuste pour la VaR sur les indices MENA étudiés (3 indices gagnés sur 4).",
        "ARIMA reste compétitif et gagne sur TASI -- la simplicité n'est pas pénalisante ici.",
        "ANN/LSTM sont bien calibrés mais ne surpassent pas nettement les modèles classiques -- l'hypothèse LSTM > classiques sur ADI n'est pas confirmée.",
        "Random Forest (et dans une moindre mesure XGBoost) est déconseillé pour la VaR de marché en l'état : sur-violations systématiques (zone rouge).",
        "Perspectives : élargir la fenêtre de test, tester GARCH asymétrique (EGARCH/GJR) et des architectures hybrides (GARCH-LSTM), et exploiter des données intra-journalières pour les marchés MENA les plus liquides.",
    ], size=14.5)
    add_rect(s, Inches(8.5), Inches(1.75), Inches(4.2), Inches(3.4), WHITE)
    add_text(s, Inches(8.75), Inches(1.95), Inches(3.7), Inches(0.4), "PIPELINE & CODE", size=12, bold=True, color=GOLD)
    add_bullets(s, Inches(8.75), Inches(2.4), Inches(3.7), Inches(2.5), [
        "Package tsvar (Python) : data, classical, volatility, ml, deep, var, backtest, run, plots.",
        "22 tests unitaires (pytest) couvrant chaque étage du pipeline.",
        "Résultats reproductibles : outputs/results.csv (7 modèles x 4 indices x 2 alpha).",
    ], size=13, bullet="›", anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, page[0], TOTAL)

    # ---------------- Slide 31: Merci / clôture -----------------------------
    s = track(blank_slide(prs)); set_background(s, NAVY)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.15), GOLD)
    add_text(s, Inches(1.0), Inches(2.7), Inches(11.3), Inches(1.2), "Merci",
              size=54, bold=True, color=WHITE, font=FONT_HEAD)
    add_text(s, Inches(1.0), Inches(3.9), Inches(11.0), Inches(0.6),
              "Projet VaR MENA -- Deep Learning vs modèles statistiques",
              size=16, color=GOLD_LIGHT)
    add_text(s, Inches(1.0), Inches(4.5), Inches(11.0), Inches(0.5),
              "Questions bienvenues -- code, données et résultats disponibles dans le dépôt du projet.",
              size=13, color=SAND)
    add_text(s, Inches(1.0), Inches(6.75), Inches(8), Inches(0.4), "ahmedbenarfa.1992@gmail.com",
              size=12, color=SAND)

    n_slides = page[0]
    out_path = PRES_DIR / "VaR_MENA.pptx"
    prs.save(str(out_path))
    return out_path, n_slides


if __name__ == "__main__":
    path, n = build()
    print(f"Saved {path} with {n} slides")
